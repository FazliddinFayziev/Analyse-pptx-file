from pptx import Presentation
from pptx.util import Pt, Inches
import requests
from io import BytesIO
import json
from pptx.dml.color import RGBColor


def set_gradient_background(slide):
    background = slide.background
    fill = background.fill
    fill.gradient()

    gradient_stops = fill.gradient_stops

    stop1 = gradient_stops[0]
    stop1.position = 0
    stop1.color.rgb = RGBColor(72, 207, 173)  # greenish

    stop2 = gradient_stops[1]
    stop2.position = 1
    stop2.color.rgb = RGBColor(79, 193, 233)  # bluvish

    fill.gradient_angle = 45  # 45 degrees
    

def emu_to_pt(value):
    return round(value / 12700)

def remove_images_from_slide(slide):
    i = 0
    while i < len(slide.shapes):
        shape = slide.shapes[i]
        if shape.shape_type == 13:  # 13 corresponds to the picture shape type
            sp = shape._element
            sp.getparent().remove(sp)
        else:
            i += 1

def replace_text_in_shapes(slide, shapes_data):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for shape_data in shapes_data:
            if shape.shape_id == shape_data.get("id"):
                text_frame = shape.text_frame
                text_frame.clear()  # Clear existing text
                p = text_frame.paragraphs[0]
                run = p.add_run()
                run.text = shape_data["text_content"]
                
                if "font_size" in shape_data:
                    # Set the font size
                    run.font.size = Pt(shape_data["font_size"])

def add_images_to_slide(slide, images_data):
    for image_data in images_data:
        image_url = image_data["path"]
        response = requests.get(image_url)

        if response.status_code == 200:
            image = BytesIO(response.content)
            left = image_data["left"]
            top = image_data["top"]
            width = image_data["width"]
            height = image_data["height"]
            slide.shapes.add_picture(image, left, top, width, height)

def replace_content_in_presentation(pptx_path, json_file_path):
    with open(json_file_path, 'r') as json_file:
        slides_data = json.load(json_file)

    prs = Presentation(pptx_path)

    for slide_data in slides_data:
        slide = prs.slides[slide_data["slide_index"] - 1]
        
        # Remove old images
        remove_images_from_slide(slide)
        # set_gradient_background(slide)
        
        # Replace text in shapes
        replace_text_in_shapes(slide, slide_data["shapes"])
        
        # Add new images
        add_images_to_slide(slide, slide_data["image_path"])

    modified_pptx_path = 'modified_presentation.pptx'
    prs.save(modified_pptx_path)
    return modified_pptx_path

# Example usage
json_file_path = 'slides_data.json'
pptx_path = replace_content_in_presentation('plan.pptx', json_file_path)
print(f'Modified presentation saved as {pptx_path}')