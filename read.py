from pptx import Presentation
import json

def get_text_boxes_info(pptx_path):
    presentation = Presentation(pptx_path)
    all_slides = []

    for slide_index, slide in enumerate(presentation.slides):
        slide_data = {
            'slide_index': slide_index + 1,
            'shapes': []
        }

        index = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_content = ""
                text_frame = shape.text_frame
                width = shape.width
                height = shape.height
                font_size = None

                for paragraph_index, paragraph in enumerate(text_frame.paragraphs):
                    for run_index, run in enumerate(paragraph.runs):
                        text_content += run.text

                        # Add newline character if it's the end of a paragraph
                        if run_index == len(paragraph.runs) - 1 and paragraph_index != len(text_frame.paragraphs) - 1:
                            text_content += "\n"

                        font_size = run.font.size
                        # Print font size if it's set
                        if font_size:
                            print(f"Font size: {font_size}")

                text_box_info = {
                    'id': index,
                    'text_content': text_content,
                    'width': width,
                    'height': height,
                    'font_size': font_size
                }

                slide_data['shapes'].append(text_box_info)
                index += 1

        all_slides.append(slide_data)

    return all_slides

# Example usage
slides_data = get_text_boxes_info('example.pptx')

# Print the resulting JSON data
print(json.dumps(slides_data, indent=2))
