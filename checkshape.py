from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# Function to get font sizes from a slide
def get_text_boxes_font_sizes(slide):
        for shape in slide.shapes:
            print(shape.shape_type)


# Example usage
pptx_path = 'example.pptx'
presentation = Presentation(pptx_path)

for slide in presentation.slides:
    get_text_boxes_font_sizes(slide)
