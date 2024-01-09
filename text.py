import json
from pptx import Presentation
from pptx.util import Inches, Pt

# Function to extract text and font size from a shape
def get_shape_details(shape):
    text_content = ""
    font_size = None
    
    if shape.has_text_frame:
        paragraphs = []
        for paragraph in shape.text_frame.paragraphs:
            runs = [run.text for run in paragraph.runs]
            paragraphs.append(' '.join(runs))
        text_content = '\n'.join(paragraphs)
        
        # Getting font size from the first run of the first paragraph
        if shape.text_frame.paragraphs[0].runs:
            font_size = shape.text_frame.paragraphs[0].runs[0].font.size
            font_size = round(font_size.pt, 2) if font_size else None
    
    return text_content, font_size

# Function to get image details
def get_image_details(shape):
    image_path = {
        "path": "",  # Placeholder for the image path, which might be extracted differently
        "width": round(shape.width * Inches(1), 2),
        "height": round(shape.height * Inches(1), 2),
        "top": round(shape.top * Inches(1), 2),
        "left": round(shape.left * Inches(1), 2)
    }
    return image_path

# Function to process each slide in the presentation
def process_slides(presentation):
    json_data = []
    for slide_index, slide in enumerate(presentation.slides):
        slide_data = {
            "slide_index": slide_index + 1,
            "shapes": [],
            "image_path": []
        }
        
        for shape_id, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # This is a picture
                slide_data["image_path"].append(get_image_details(shape))
            else:
                text_content, font_size = get_shape_details(shape)
                shape_data = {
                    "id": shape_id,
                    "width": round(shape.width * Inches(1), 2),
                    "height": round(shape.height * Inches(1), 2),
                    "text_content": text_content,
                    "font_size": font_size
                }
                slide_data["shapes"].append(shape_data)
        
        json_data.append(slide_data)
    return json_data

def pptx_to_json(pptx_file_path):
    presentation = Presentation(pptx_file_path)
    json_data = process_slides(presentation)
    return json.dumps(json_data, indent=2)

# Specify the path to your pptx file
pptx_file_path = 'plan.pptx'

# Convert the PPTX to JSON
json_output = pptx_to_json(pptx_file_path)

# Output the JSON to a file
with open('output.json', 'w') as json_file:
    json_file.write(json_output)

print("JSON data has been generated and written to output.json")