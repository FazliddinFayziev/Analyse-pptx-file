from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

def create_combined_charts(products, bar_data, doughnut_data, 
                           bar_chart_title='Bar Chart', doughnut_chart_title='Doughnut Chart', 
                           output_filename='CombinedCharts.pptx', positions=None, show_values=True, years=None):
    presentation = Presentation()
    
    total_charts = max(len(bar_data), len(doughnut_data))
    charts_per_slide = len(positions)
    slides_required = -(-total_charts // charts_per_slide)

    for slide_num in range(slides_required):
        slide_layout = presentation.slide_layouts[6] 
        slide = presentation.slides.add_slide(slide_layout)

        num_charts_on_slide = min(charts_per_slide, total_charts - slide_num * charts_per_slide)

        for i in range(num_charts_on_slide):
            x, y, width, height = positions[i]

            # Determine chart type for the current position
            if i < len(bar_data):
                chart_type = XL_CHART_TYPE.BAR_CLUSTERED
                data = bar_data[i]
                chart_title = f'{bar_chart_title} - {years[i]}' if years else bar_chart_title
            else:
                chart_type = XL_CHART_TYPE.DOUGHNUT
                data = doughnut_data[i - len(bar_data)]
                chart_title = f'{doughnut_chart_title} - {years[i]}' if years else doughnut_chart_title

            # Add chart to the slide
            chart_data = CategoryChartData()
            chart_data.categories = products
            chart_data.add_series(f'Series {i + 1}', data)

            chart = slide.shapes.add_chart(
                chart_type=chart_type, x=x, y=y, cx=width, cy=height, chart_data=chart_data
            ).chart

            chart.has_title = True
            chart.chart_title.text_frame.text = chart_title

            if show_values:
                for j, point in enumerate(chart.series[0].points):
                    if j < len(data):
                        value = point.data_label.text_frame
                        value.text = f'{data[j]}%'
                        value.paragraphs[0].font.size = Pt(8)

    presentation.save(output_filename)

    print(f"Combined charts created and saved in '{output_filename}'")

# Example usage
products = ['Product A', 'Product B', 'Product C', 'Product D', 'Product F']
bar_data_single = [10, 20, 30, 10, 5, 26]
bar_data_multiple = [
    [15, 25, 35, 25],
    [5, 15, 25, 15],
    [10, 20, 30, 10],
    [8, 18, 28, 8],
    [12, 22, 32, 12],
    [12, 22, 32, 12],
]

doughnut_data_single = [10, 20, 30, 10, 5, 26]
doughnut_data_multiple = [
    [15, 25, 35, 25],
    [5, 15, 25, 15],
    [10, 20, 30, 10],
    [8, 18, 28, 8],
    [12, 22, 32, 12],
    [12, 22, 32, 12],
]

years = [2021, 2022, 2023, 2024, 2025, 2026]

# Define custom positions for charts
positions_custom = [
    (Inches(0.4), Inches(0.2), Inches(4), Inches(3)),
    (Inches(4.6), Inches(0.2), Inches(4), Inches(3)),
    (Inches(0.4), Inches(3.7), Inches(4), Inches(3)),
    (Inches(4.6), Inches(3.7), Inches(4), Inches(3)),
]

# Combine bar and doughnut charts with custom positions
create_combined_charts(products, bar_data_multiple, doughnut_data_multiple, 
                       bar_chart_title='Bar Chart', doughnut_chart_title='Doughnut Chart',
                       positions=positions_custom, show_values=True, years=years, 
                       output_filename='CombinedCharts_CustomPosition.pptx')
