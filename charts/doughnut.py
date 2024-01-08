from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt

def create_doughnut_chart(products, data, chart_title='Product Distribution', output_filename='DoughnutChart.pptx', positions=None, show_values=True, years=None):
    presentation = Presentation()
    total_charts = len(data)
    charts_per_slide = len(positions)  
    slides_required = -(-total_charts // charts_per_slide) 

    for slide_num in range(slides_required):
        slide_layout = presentation.slide_layouts[6]
        slide = presentation.slides.add_slide(slide_layout)

        num_charts_on_slide = min(charts_per_slide, total_charts - slide_num * charts_per_slide)

        for i in range(num_charts_on_slide):
            x, y, width, height = positions[i]

            chart_data = CategoryChartData()
            chart_data.categories = products
            chart_data.add_series(f'Series {i + 1}', data[slide_num * charts_per_slide + i])

            chart = slide.shapes.add_chart(
                chart_type=XL_CHART_TYPE.DOUGHNUT, x=x, y=y, cx=width, cy=height, chart_data=chart_data
            ).chart

            chart.has_title = True
            chart.chart_title.text_frame.text = f'{chart_title} - {years[slide_num * charts_per_slide + i]}' if years else chart_title

            if show_values:
                for j, point in enumerate(chart.series[0].points):
                    if j < len(data[slide_num * charts_per_slide + i]):
                        value = point.data_label.text_frame
                        value.text = f'{data[slide_num * charts_per_slide + i][j]}%'
                        value.paragraphs[0].font.size = Pt(8)
                        
    presentation.save(output_filename)

    print(f"Doughnut chart(s) created and saved in '{output_filename}'")

# Example usage
products = ['Product A', 'Product B', 'Product C', 'Product D']
# data_single = [10, 20, 30, 10, 5, 26]
data_multiple = [
    [15, 25, 35, 25],
    # [5, 15, 25, 15],
    # [10, 20, 30, 10],
    # [8, 18, 28, 8],
    # [12, 22, 32, 12],
    # [12, 22, 32, 12],
]
years = [2021]

# Define custom positions for charts
positions_custom = [
    (Inches(0.4), Inches(0.2), Inches(4), Inches(3)),
    # (Inches(4.6), Inches(0.2), Inches(4), Inches(3)),
    # (Inches(0.4), Inches(3.7), Inches(4), Inches(3)),
    # (Inches(4.6), Inches(3.7), Inches(4), Inches(3)),
]

# Single data series and no years with custom positions
# create_doughnut_chart(products, [data_single], chart_title='Single Chart', positions=positions_custom, show_values=True, output_filename='SingleDoughnutChart_CustomPosition.pptx')

# Multiple data series and years with custom positions
create_doughnut_chart(products, data_multiple, chart_title='Multiple Charts', positions=positions_custom, show_values=True, years=years, output_filename='done.pptx')
