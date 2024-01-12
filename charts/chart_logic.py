from charts import charts
from data import all_data
from pptx.util import Inches
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from functions import calculating_all_positions

# create presentation
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
def create_charts(slide, all_charts):
    W = 10
    H = 7.5
    
    for j, chart_info in enumerate(all_charts):
        w = chart_info["size"][0]
        h = chart_info["size"][1]
        position = chart_info["positions"][0]
        product_names = chart_info["productName"]
        chart_type = charts[chart_info["chartType"]]
        
        if type(chart_info["positions"][0]) == str:
            [left, top] = calculating_all_positions(position, W, H, w, h)
        else:
            [left, top] = chart_info["positions"]
        
        chart_data = CategoryChartData()
        chart_data.categories = [category for category in chart_info["categories"]]
        
        for product, series_data in zip(product_names, chart_info["data"]):
            chart_data.add_series(product, series_data)
            
        chart = slide.shapes.add_chart(
            chart_type, Inches(left), Inches(top), Inches(w), Inches(h), chart_data
        ).chart
        
        
        if "show" in chart_info:
            chart.plots[0].has_data_labels = chart_info["show"]
        else:
            chart.plots[0].has_data_labels = False
        
        # Add title to the chart
        chart.has_title = True
        chart.chart_title.text_frame.text = chart_info["chartTitle"]
        
        if "legend" in chart_info:
            chart.has_legend = chart_info["legend"]
        else:
            chart.has_legend = True

# Example data for a line chart with two series
line_chart_data = [
    {
        "chartType": "RADAR_MARKERS",
        "productName": ["Series 1", "Series 2"],
        "data": [[10, 15, 20]],
        "chartTitle": 'Line Chart Example',
        "size": [8, 6],
        # Custom position [4, 3]
        "positions": ["CENTER"],
        "categories": [2001],
        # "show": True,
        "legend": True,
    }
]

create_charts(slide, all_data)
create_charts(slide2, line_chart_data)
prs.save('done.pptx')
print('It is done')
