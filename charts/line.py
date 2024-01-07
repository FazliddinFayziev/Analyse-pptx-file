from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

def create_line_chart(products, data, years, chart_title='Product Distribution', 
                      output_filename='LineChart.pptx', width=4, height=3, show_values=True, position='center'):
    # Create PowerPoint presentation
    presentation = Presentation()

    # Add a slide
    slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Calculate position based on the specified option
    if position == 'center':
        x = Inches((10 - width) / 2)
        y = Inches((7.5 - height) / 2)
    elif position == 'top-left':
        x = Inches(1)
        y = Inches(1)
    else:
        raise ValueError("Invalid position option. Choose 'center' or 'top-left'.")

    cx, cy = Inches(width), Inches(height)

    # Add Line chart to the slide
    chart_data = CategoryChartData()
    chart_data.categories = years

    for i, series_data in enumerate(data):
        chart_data.add_series(f'{products[i]}', series_data)

    chart = slide.shapes.add_chart(
        chart_type=XL_CHART_TYPE.LINE, x=x, y=y, cx=cx, cy=cy, chart_data=chart_data
    ).chart

    # Set the chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title

    # Show values on data points
    if show_values:
        for series in chart.series:
            for i, point in enumerate(series.points):
                value = point.data_label.text_frame
                value.text = f'{data[series.index][i]}%'
                value.paragraphs[0].font.size = Inches(0.2)

    # Set period between data points
    chart.series[0].data_labels.show_category_name = False
    chart.series[0].data_labels.show_value = False
    chart.series[0].data_labels.show_legend_key = False
    chart.series[0].data_labels.show_percentage = False
    chart.series[0].data_labels.show_leader_lines = False

    # Center and justify the chart on the slide
    chart.left = x
    chart.top = y

    # Save the PowerPoint presentation
    presentation.save(output_filename)

    print(f"Line chart created and saved in '{output_filename}'")

# Example usage with multiple data series and years
products = ['Product A', 'Product B', 'Product C', 'Product D']
data = [
    [10, 20, 30, 10],
    [5, 15, 25, 15],
    [15, 25, 35, 25]
]
years = [2021, 2022, 2023, 2024]

create_line_chart(products, data, years, width=8, height=6, show_values=True, position='center')
