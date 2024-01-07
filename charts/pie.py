from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor

def create_pie_chart(products, data, colors, chart_title='Product Distribution', 
                     output_filename='PieChart.pptx', width=4, height=3, show_values=True, position='center'):
    # Create PowerPoint presentation
    presentation = Presentation()

    # Add a slide
    slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Calculate position based on the specified option
    if position == 'center':
        x = Inches((10 - width) / 2)
        y = Inches((7.5 - height) / 2)
    elif position == 'top-left':
        x = Inches(1)
        y = Inches(1)
    else:
        raise ValueError("Invalid position option. Choose 'center' or 'top-left'.")

    cx, cy = Inches(width), Inches(height)

    # Add Pie chart to the slide
    chart_data = CategoryChartData()
    chart_data.categories = products
    chart_data.add_series('Series 1', data)

    chart = slide.shapes.add_chart(
        chart_type=XL_CHART_TYPE.PIE, x=x, y=y, cx=cx, cy=cy, chart_data=chart_data
    ).chart

    # Set the chart title
    chart.has_title = True
    chart.chart_title.text_frame.text = chart_title

    # Set colors
    for i, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = RGBColor(*colors[i])  # Convert color to RGBColor

    # Show values within the chart
    if show_values:
        for i, point in enumerate(chart.series[0].points):
            value = point.data_label.text_frame
            value.text = f'{data[i]}%'
            value.paragraphs[0].font.size = Inches(0.2)

    # Center and justify the chart on the slide
    chart.left = x
    chart.top = y

    # Save the PowerPoint presentation
    presentation.save(output_filename)

    print(f"Pie chart created and saved in '{output_filename}'")

# Example usage
products = ['Product A', 'Product B', 'Product C', 'Product D']
data = [10, 20, 30, 10]
colors = [(102, 179, 255), (153, 255, 153), (255, 204, 153), (255, 0, 0)]  # RGB values for the colors

create_pie_chart(products, data, colors, width=8, height=6, show_values=True, position='center')
