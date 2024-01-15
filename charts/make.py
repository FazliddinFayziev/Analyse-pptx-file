from charts import charts
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_MARKER_STYLE 
from functions import position, size_value, axis, label_value, get_layout_dimensions, format_data_for_xy_scatter, chart_colors, legend_value, years, font, title

class ChartMaker:
    def make_charts(self, slide, all_charts):
        W, H = get_layout_dimensions(slide)

        for chart_info in all_charts:

            product_names = chart_info["productName"]
            chart_type = charts[chart_info["chartType"].upper()]

            # SIZE
            [w, h] = size_value(chart_info["size"])

            # POSITION
            [left, top] = position(chart_info["position"], W, H, w, h)

            # CHART
            chart_data = chart_type[1]
            
            # YEARS
            years(chart_info, chart_data)
            
            if chart_type[2] == 1:
                for product, series_data in zip(product_names, chart_info["data"]):
                    chart_data.add_series(product, series_data)
                
                chart = slide.shapes.add_chart(
                    chart_type[0], Inches(left), Inches(top), Inches(w), Inches(h), chart_data
                ).chart
                
                
            elif chart_type[2] == 2:
                for product, series_data in zip(product_names, chart_info["data"]):
                    series = chart_data.add_series(product)
                    x_value, y_value, size = series_data
                    series.add_data_point(x_value, y_value, size)

                chart = slide.shapes.add_chart(
                    chart_type[0], Inches(left), Inches(top), Inches(w), Inches(h), chart_data
                ).chart
                
            
            elif chart_type[2] == 3:
                data = format_data_for_xy_scatter(chart_info["data"], chart_info["years"])
                for product, series_data in zip(product_names, data):
                    series = chart_data.add_series(product)
                    for point in series_data:
                        series.add_data_point(point[0], point[1])

                chart = slide.shapes.add_chart(
                    chart_type[0], Inches(left), Inches(top), Inches(w), Inches(h), chart_data
                ).chart

                for idx, color in enumerate(chart_info["colors"]):
                    series = chart.series[idx]
                    try:
                        series.marker.style = getattr(XL_MARKER_STYLE, chart_info["marker"][idx].upper())
                    except IndexError:
                        series.marker.style = XL_MARKER_STYLE.CIRCLE
                    series.marker.size = 8
                    fill = series.format.fill
                    fill.solid()
                    fill.fore_color.rgb = RGBColor.from_string(color[1:])

            # COLORS
            chart_colors(chart_info, chart)

            # SHOW VALUE
            label_value(chart_info, chart)

            # TITLE
            title(chart, chart_info["chartTitle"])

            # FONT LOGIC
            font(chart)

            # LEGEND
            legend_value(chart_info, chart)
            
            # AXIS
            # axis(chart, chart_info)
            
                