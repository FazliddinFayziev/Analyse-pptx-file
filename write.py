from pptx import Presentation
from pptx.util import Pt
import requests
from io import BytesIO
import json

def emu_to_pt(value):
    return round(value / 12700)

def area(width, height):
    w = emu_to_pt(width)
    h = emu_to_pt(height)
    return w * h

def text_covered(text, font_size):
    return text * (font_size ** 2)

def reduce_font_size(font_size):
    return font_size - 1

def replace_text_in_text_frames(pptx_path, json_file_path):
    with open(json_file_path, 'r') as json_file:
        new_texts = json.load(json_file)

    prs = Presentation(pptx_path)

    for new_text in new_texts:
        index_slide = new_text["slide_index"] - 1
        slide = prs.slides[index_slide]
        text_index = 0
        
        # Remove old images
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                slide.shapes._spTree.remove(shape._element)

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
      
            if text_index < len(new_text["shapes"]):
                text_content_data = new_text["shapes"][text_index]
                run = p.add_run()
                run.text = text_content_data["text_content"]
        
                if "font_size" in text_content_data:
                    font_size_value = text_content_data["font_size"]
                    text = text_covered(len(run.text), font_size_value)
                    square = area(text_content_data["width"], text_content_data["height"])
        
                    while text > square:
                        font_size_value = reduce_font_size(font_size_value)
                        text = text_covered(len(run.text), font_size_value)
                        square = area(text_content_data["width"], text_content_data["height"])
                    
                    run.font.size = Pt(font_size_value)
                text_index += 1
        
            if new_text["image_path"]:
                for image_index, image in enumerate(new_text["image_path"]):
                    image_url = image["path"]
                    response = requests.get(image_url)

                    if response.status_code == 200:
                        image_data = BytesIO(response.content)
                        left = image["left"]
                        top = image["top"]
                        width = image["width"]
                        height = image["height"]
                        slide.shapes.add_picture(image_data, left, top, width, height)

                new_text["image_path"] = []
                text_index += 1
              
    modified_pptx_path = 'modified_education_presentation_fixed.pptx'
    prs.save(modified_pptx_path)
    return modified_pptx_path

# Example usage
json_file_path = 'slides_data.json'
replace_text_in_text_frames('example.pptx', json_file_path)
print('Done')
