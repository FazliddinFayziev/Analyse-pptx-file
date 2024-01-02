from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Creating presentation object
root = Presentation('edu_temp.pptx')

# Access the first slide in the presentation
slide = root.slides[0]

titles_and_sizes = [{'text': 'Elegant Education Pack for Students ', 'font_size': 32, 'font_color': 'Automatic', 'font_family': 'Arial', 'text_frame_type': "title", 'left': 1283100, 'top': 458100, 'width': 6577800, 'height': 2571750},
                    {'text': 'Here is where your presentation begins', 'font_size': 28, 'font_color': 'Automatic', 'font_family': 'Calibri', 'text_frame_type': "content", 'left': 1283100, 'top': 3394625, 'width': 6577800, 'height': 458100},
                    {'text': 'Done by Fazliddin', 'font_size': 24, 'font_color': 'Automatic', 'font_family': 'Times New Roman', 'text_frame_type': "content", 'left': -1377443, 'top': 4685400, 'width': 6577800, 'height': 458100}]

for data in titles_and_sizes:
    # Create a new text frame for each text on the existing slide
    if data["text_frame_type"] == "title":
        if slide.shapes.title.text_frame is not None:
            text_frame = slide.shapes.title.text_frame
            text_frame.clear()  # Clear any existing text
        else:
            text_frame = slide.shapes.title.text_frame
    else:
        text_frame = slide.shapes.add_textbox(data["left"], data["top"], data["width"], data["height"]).text_frame

    # Set the autofit property to ResizeTextOnOverflow
    text_frame.word_wrap = True
    text_frame.auto_size = True

    # Add a new paragraph for each text
    paragraph = text_frame.add_paragraph()

    # Set the text content, font size, and font family for each paragraph
    paragraph.text = data["text"]
    paragraph.font.size = Pt(data["font_size"])
    paragraph.font.color.rgb = data["font_color"]
    paragraph.font.name = data["font_family"]

    # Set paragraph alignment to center
    paragraph.alignment = PP_ALIGN.CENTER

    # Set text frame alignment to center
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Save the presentation
root.save("Output.pptx")

print("done")