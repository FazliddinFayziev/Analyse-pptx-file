from pptx import Presentation
from pptx.util import Pt

def replace_text_with_auto_formatting(paragraph, old_text, new_text, font_size=None, bold=False, italic=False, underline=False):
    for run in paragraph.runs:
        if old_text in run.text:
            run.text = new_text
            run.font.size = font_size if font_size else run.font.size
            run.font.bold = bold
            run.font.italic = italic
            run.font.underline = underline

def adjust_font_size(paragraph, new_text, max_width, default_font_size=Pt(18), min_font_size=Pt(8)):
    run = paragraph.runs[0]
    run.text = new_text
    run.font.size = default_font_size

    while paragraph.width > max_width and run.font.size > min_font_size:
        run.font.size -= Pt(1)

def replace_text_in_slide(slide, replacements):
    for replacement in replacements:
        text_frame = None
        if "shape_index" in replacement and replacement["shape_index"] < len(slide.shapes):
            shape = slide.shapes[replacement["shape_index"]]
            if shape.has_text_frame:
                text_frame = shape.text_frame
        elif "shape_name" in replacement:
            shape = slide.shapes.get(replacement["shape_name"])
            if shape and shape.has_text_frame:
                text_frame = shape.text_frame

        if text_frame:
            for i, text_content in enumerate(replacement["text_contents"]):
                if i < len(text_frame.paragraphs):
                    paragraph = text_frame.paragraphs[i]
                    replace_text_with_auto_formatting(paragraph, text_content["old_text"], text_content["new_text"], font_size=text_content.get("font_size"))

                    # Check if text overflows the paragraph and adjust font size
                    max_width = text_content.get("max_width", text_frame.width)
                    if paragraph.width > max_width and text_content.get("adjust_font_size", True):
                        adjust_font_size(paragraph, text_content["new_text"], max_width)

def main():
    presentation_path = "path/to/your/presentation.pptx"
    presentation = Presentation(presentation_path)

    replacements = [
        {
            "slide_index": 1, "shape_index": 0, "text_contents": [
                {"old_text": "Old Text 1", "new_text": "Elegant Education Pack for Students in Malaysia and United States", "max_width": 6577800},
                {"old_text": "Old Text 2", "new_text": "Here is where your presentation begins", "max_width": 6577800},
                {"old_text": "Old Text 3", "new_text": "Done by Fazliddin", "max_width": 6577800}
            ]
        }
        # Add more replacements as needed
    ]

    for replacement in replacements:
        if replacement["slide_index"] - 1 < len(presentation.slides):
            slide = presentation.slides[replacement["slide_index"] - 1]
            replace_text_in_slide(slide, [replacement])

    output_path = "education.pptx"
    presentation.save(output_path)

if __name__ == "__main__":
    main()
